from PyQt5 import Qt, QtWidgets, QtGui, QtCore
from PyQt5.QtCore import QSize
from PyQt5.QtWidgets import QPushButton
import sys
from PyQt5.QtGui import QPainter, QColor, QIcon
from pptx import Presentation
from fpdf import FPDF
from PyQt5.QtWidgets import QFileDialog
from PyQt5.QtGui import QPixmap
from PIL import Image


class Slide:
    def __init__(self, type, title='', text='', picture=None):
        self.type = type
        self.text = text
        self.title = title
        self.picture = picture

    def set_text(self, text):
        self.text = text

    def set_title(self, text):
        self.title = text

    def set_picture(self, picture):
        self.picture = picture


class GreetingWindow(QtWidgets.QMainWindow):
    resized = QtCore.pyqtSignal()

    def __init__(self):
        super().__init__()
        self.initUI()

    def initUI(self):
        self.resized.connect(self.resizing)
        self.resize(1100, 900)
        self.setWindowTitle('Презентация')

        self.pushButton = QPushButton(self)
        self.pushButton.setText('Создать презентацию')
        font = QtGui.QFont()
        font.setFamily('comic sans ms')
        font.setPointSize(33)
        self.pushButton.setFont(font)
        self.pushButton.setStyleSheet(('''QPushButton {            
                                    background: rgb(184,232,176);
                                    position: absolute;
                                    border-radius: 70px;
                                    padding: 10px;
                                    color: rgb(255, 255, 255);
                                }'''))
        self.pushButton.clicked.connect(self.create_presentation)

        self.label = QtWidgets.QLabel(self)
        self.label.setText('Добро пожаловать!')
        font = QtGui.QFont()
        font.setPointSize(200)
        self.label.setFont(QtGui.QFont("comic sans ms", 46, QtGui.QFont.Bold))
        self.label.setStyleSheet("color: rgb(184,232,176);")
        self.label.adjustSize()

    def create_presentation(self):
        self.k = ChooseTemplate([], (self.width(), self.height()))
        self.k.show()
        self.hide()

    def resizing(self):
        self.pushButton.setGeometry(int(0.1606 * self.width()),
                                    int(0.5000 * self.height()),
                                    int(0.6569 * self.width()),
                                    int(0.3900 * self.height()))
        self.label.setGeometry(int(0.193 * self.width()),
                               int(0.065 * self.height()),
                               int(0.7 * self.width()),
                               int(0.4 * self.height()))

    def resizeEvent(self, event):
        self.resized.emit()
        return super().resizeEvent(event)


class MainScreen(QtWidgets.QMainWindow):
    resized = QtCore.pyqtSignal()

    def __init__(self, slides, size):
        self.size_of_window = size
        super().__init__()
        self.slides = slides
        self.number_of_slides = len(self.slides)
        self.initUI()

    def initUI(self):
        self.resized.connect(self.resizing)
        self.resize(*self.size_of_window)
        self.setWindowTitle('Презентация')

        self.pushButton = QPushButton('Готово', self)
        self.pushButton.setStyleSheet('''QPushButton {            
                                    background: rgb(13,174,78);
                                    border-style: outset;
                                    border-width: px;
                                    border-radius: 15px;
                                    border-color: rgb(49,106,60);
                                    padding: 4px;
                                    color: rgb(255, 255, 255);
                                }''')
        font = QtGui.QFont()
        font.setPointSize(12)
        self.pushButton.setFont(font)
        self.pushButton.show()

        layout = Qt.QGridLayout(self)

        self.list_of_buttons = []

        for i in range(self.number_of_slides):
            x = QtWidgets.QPushButton(self)
            x.setGeometry(int(0.0089 * self.width()), int(0.078 * i * self.height()), int(0.1387 * self.width()), int(0.1078 * self.height()))
            x.setStyleSheet('''QPushButton {
                                    background: rgb(95,193,215);
                                    border-style: outset;
                                    border-width: 2px;
                                    border-radius: 25px;
                                    border-color: white ;
                                    padding: 4px;
                                    color: rgb(255, 255, 255);
                                }''')
            font = QtGui.QFont()
            font.setPointSize(14)
            x.setFont(font)
            x.setMaximumSize(int(0.1387 * self.width()), int(0.1450 * self.height()))
            x.setMinimumSize(int(0.1387 * self.width()), int(0.1450 * self.height()))
            x.setText(self.slides[i].title)
            self.list_of_buttons.append(x)
            layout.addWidget(x, i, 0)

        self.create_new_slide = QtWidgets.QPushButton(self)
        self.create_new_slide.setGeometry(int(0.0089 * self.width()), int(0.078 * self.number_of_slides * self.height()), int(0.1387 * self.width()), int(0.1078 * self.height()))
        self.create_new_slide.setStyleSheet('''QPushButton {      
                                    background: rgb(95,193,215);
                                    border-style: outset;
                                    border-width: 2px;
                                    border-radius: 25px;
                                    padding: 4px;
                                    color: rgb(255, 255, 255);
                                }''')  #border-style: outset;border-width: 2px;border-radius: 25px;padding: 4px;(параметры для границы(обводки))

        font = QtGui.QFont()
        font.setPointSize(50)
        self.create_new_slide.setFont(font)
        self.create_new_slide.setMaximumSize(int(0.1387 * self.width()), int(0.1450 * self.height()))
        self.create_new_slide.setMinimumSize(int(0.1387 * self.width()), int(0.1450 * self.height()))
        self.create_new_slide.setText('+')
        layout.addWidget(self.create_new_slide, self.number_of_slides, 0)

        w = Qt.QWidget(self)
        w.setLayout(layout)

        self.mw = Qt.QScrollArea(self)
        self.mw.setWidget(w)
        self.mw.setStyleSheet('background: rgb(95,193,215);')
        self.mw.setGeometry(0, 0, int(0.1825 * self.width()), int(0.9196 * self.height()))

    def paintEvent(self, event):
        self.qp = QPainter()
        self.qp.begin(self)
        self.draw(event, self.qp)
        self.qp.end()

    def draw(self, event, qp):
        qp.setPen(QColor(0, 0, 0))
        qp.drawRect(0, 0, int(0.1825 * self.width()), int(1 * self.height()))

    def resizing(self):
        self.pushButton.setGeometry(int(0.0250 * self.width()),
                                    int(0.9374 * self.height()),
                                    int(0.1000 * self.width()),
                                    int(0.0500 * self.height()),
                                    )
        self.mw.resize(int(0.1825 * self.width()), int(0.9196 * self.height()))

    def resizeEvent(self, event):
        self.resized.emit()
        return super().resizeEvent(event)

    def create_presentation(self):
        filename, ok = QFileDialog.getSaveFileName(self, "Сохранить файл", ".", "Презентация(*.pptx);;PDF(*.pdf)")
        if filename.endswith('pptx'):
            prs = Presentation()
            slide1 = prs.slide_layouts[0]
            slide2 = prs.slide_layouts[1]
            slide3 = prs.slide_layouts[5]
            slide4 = prs.slide_layouts[3]
            for i in self.slides:
                if i.type == 1:
                    slide = prs.slides.add_slide(slide1)
                    slide.shapes.title.text = i.title
                    slide.placeholders[1].text = i.text
                elif i.type == 2:
                    slide = prs.slides.add_slide(slide2)
                    slide.shapes.title.text = i.title
                    slide.placeholders[1].text = i.text
                elif i.type == 3:
                    slide = prs.slides.add_slide(slide3)
                    slide.shapes.title.text = i.title
                    if i.picture:
                        im = Image.open(i.picture)
                        if im.size[0] / im.size[1] > 80/48:
                            picture = slide.shapes.add_picture(i.picture, 500000, 1600000, 8000000)
                        else:
                            picture = slide.shapes.add_picture(i.picture, 500000, 1600000, height=4800000)
                else:
                    slide = prs.slides.add_slide(slide4)
                    slide.shapes.title.text = i.title
                    slide.placeholders[2].text = i.text
                    if i.picture:
                        im = Image.open(i.picture)
                        if im.size[0] / im.size[1] > 40 / 48:
                            picture = slide.shapes.add_picture(i.picture, 500000, 1600000, 4000000)
                        else:
                            picture = slide.shapes.add_picture(i.picture, 500000, 1600000, height=4800000)
            prs.save(filename)
        else:
            prs = FPDF(orientation='L', unit='mm', format='A4')
            for i in self.slides:
                prs.add_page()
                prs.add_font('DejaVu', '', 'DejaVuSansCondensed.ttf', uni=True)
                prs.set_font('DejaVu', '', 48)
                prs.cell(297, 25, txt=i.title, ln=1, align="C")
                if i.type == 1:
                    prs.set_font('DejaVu', '', size=24)
                    prs.cell(280, 150, txt=i.text, ln=1, align="C")
                if i.type == 2:
                    prs.set_font('DejaVu', '', size=24)
                    prs.cell(280, 100, txt=i.text, ln=1, align="L")
                if i.type == 3:
                    if i.picture:
                        im = Image.open(i.picture)
                        if im.size[0] / im.size[1] > 26 / 17:
                            prs.image(i.picture, 10, 35, w=260)
                        else:
                            prs.image(i.picture, 10, 35, h=170)
                if i.type == 4:
                    prs.set_font('DejaVu', '', size=24)
                    prs.x += 155
                    prs.cell(140, 100, txt=i.text, ln=1, align="L")
                    prs.x += 155
                    if i.picture:
                        im = Image.open(i.picture)
                        if im.size[0] / im.size[1] > 13 / 17:
                            prs.image(i.picture, 10, 35, w=130)
                        else:
                            prs.image(i.picture, 10, 35, h=170)
            prs.output(filename)

    def to_slide(self):
        sender = self.sender()
        i = self.list_of_buttons.index(sender)
        slide = self.slides[i]
        if slide.type == 1:
            self.a = Template1(self.slides, i, (self.width(), self.height()))
        elif slide.type == 2:
            self.a = Template2(self.slides, i, (self.width(), self.height()))
        elif slide.type == 3:
            self.a = Template3(self.slides, i, (self.width(), self.height()))
        else:
            self.a = Template4(self.slides, i, (self.width(), self.height()))
        self.a.show()
        self.hide()

    def delete_slide(self):
        del self.slides[self.number_of_slide]
        if self.number_of_slide < len(self.slides):
            if self.slides[self.number_of_slide].type == 1:
                self.a = Template1(self.slides, self.number_of_slide, (self.width(), self.height()))
                self.a.show()
                self.hide()
            if self.slides[self.number_of_slide].type == 2:
                self.a = Template2(self.slides, self.number_of_slide, (self.width(), self.height()))
                self.a.show()
                self.hide()
            if self.slides[self.number_of_slide].type == 3:
                self.a = Template3(self.slides, self.number_of_slide, (self.width(), self.height()))
                self.a.show()
                self.hide()
            if self.slides[self.number_of_slide].type == 4:
                self.a = Template4(self.slides, self.number_of_slide, (self.width(), self.height()))
                self.a.show()
                self.hide()
        elif len(self.slides) >= 1:
            self.number_of_slide -= 1
            if self.slides[self.number_of_slide].type == 1:
                self.a = Template1(self.slides, self.number_of_slide, (self.width(), self.height()))
                self.a.show()
                self.hide()
            if self.slides[self.number_of_slide].type == 2:
                self.a = Template2(self.slides, self.number_of_slide, (self.width(), self.height()))
                self.a.show()
                self.hide()
            if self.slides[self.number_of_slide].type == 3:
                self.a = Template3(self.slides, self.number_of_slide, (self.width(), self.height()))
                self.a.show()
                self.hide()
            if self.slides[self.number_of_slide].type == 4:
                self.a = Template4(self.slides, self.number_of_slide, (self.width(), self.height()))
                self.a.show()
                self.hide()
        else:
            self.a = ChooseTemplate([], (self.width(), self.height()))
            self.a.show()
            self.hide()


class ChooseTemplate(MainScreen):
    resized = QtCore.pyqtSignal()

    def __init__(self, slides, size):
        super().__init__(slides, size)
        self.tempinit()

    def tempinit(self):
        self.resized.connect(self.resizing1)
        for i in self.list_of_buttons:
            i.clicked.connect(self.to_slide)

        self.label = QtWidgets.QLabel(self)
        self.label.setText('Выберите макет слайда')
        font = QtGui.QFont()
        font.setPointSize(60)
        self.label.setFont(QtGui.QFont("comic sans ms", 30, QtGui.QFont.Bold))
        self.label.setStyleSheet("color: rgb(184,232,176);")
        self.label.adjustSize()

        self.tmp1 = QPushButton("", self)
        self.tmp1.setIcon(QIcon('Shablon1.jpg'))
        self.tmp1.clicked.connect(self.to_template1)

        self.tmp2 = QPushButton("", self)
        self.tmp2.setIcon(QIcon('Shablon2.jpg'))
        self.tmp2.clicked.connect(self.to_template2)

        self.tmp3 = QPushButton("", self)
        self.tmp3.setIcon(QIcon('Shablon3.jpg'))
        self.tmp3.clicked.connect(self.to_template3)

        self.tmp4 = QPushButton("", self)
        self.tmp4.setIcon(QIcon('Shablon4.jpg'))
        self.tmp4.clicked.connect(self.to_template4)

        self.pushButton.clicked.connect(self.create_presentation)

    def to_template1(self):
        self.slides.append(Slide(1))
        self.a = Template1(self.slides, len(self.slides) - 1, (self.width(), self.height()))
        self.a.show()
        self.hide()

    def to_template2(self):
        self.slides.append(Slide(2))
        self.a = Template2(self.slides, len(self.slides) - 1, (self.width(), self.height()))
        self.a.show()
        self.hide()

    def to_template3(self):
        self.slides.append(Slide(3))
        self.a = Template3(self.slides, len(self.slides) - 1, (self.width(), self.height()))
        self.a.show()
        self.hide()

    def to_template4(self):
        self.slides.append(Slide(4))
        self.a = Template4(self.slides, len(self.slides) - 1, (self.width(), self.height()))
        self.a.show()
        self.hide()

    def resizing1(self):
        self.label.setGeometry(int(0.3410 * self.width()),
                               int(0.0714 * self.height()),
                               int(0.6963 * self.width()),
                               int(0.1100 * self.height()))
        self.tmp1.setGeometry(int(0.2670 * self.width()),
                              int(0.2150 * self.height()),
                              int(0.2600 * self.width()),
                              int(0.3500 * self.height()))
        self.tmp2.setGeometry(int(0.6400 * self.width()),
                              int(0.2150 * self.height()),
                              int(0.2600 * self.width()),
                              int(0.3500 * self.height()))
        self.tmp3.setGeometry(int(0.2670 * self.width()),
                              int(0.6072 * self.height()),
                              int(0.2600* self.width()),
                              int(0.3500 * self.height()))
        self.tmp4.setGeometry(int(0.6400 * self.width()),
                              int(0.6072 * self.height()),
                              int(0.2600* self.width()),
                              int(0.3500 * self.height()))
        self.tmp1.setIconSize(QSize(int(0.365 * self.width()), int(0.3304 * self.height())))
        self.tmp2.setIconSize(QSize(int(0.365 * self.width()), int(0.3304 * self.height())))
        self.tmp3.setIconSize(QSize(int(0.365 * self.width()), int(0.3304 * self.height())))
        self.tmp4.setIconSize(QSize(int(0.365 * self.width()), int(0.3304 * self.height())))

    def resizeEvent(self, event):
        self.resized.emit()
        return super().resizeEvent(event)


class Template1(MainScreen):
    resized = QtCore.pyqtSignal()
    def __init__(self, slides, number_of_slide, size):
        super().__init__(slides, size)
        self.number_of_slide = number_of_slide
        self.tempinit(number_of_slide)
        for i in self.list_of_buttons:
            i.clicked.connect(self.to_slide)

    def tempinit(self, number_of_slide):
        self.resized.connect(self.resizing2)
        self.print_title = QtWidgets.QPlainTextEdit(self.slides[number_of_slide].title, self)
        self.print_title.setGeometry(300, 150, 700, 250)
        font = QtGui.QFont()
        font.setPointSize(40)
        self.print_title.setFont(font)

        self.print_text = QtWidgets.QPlainTextEdit(self.slides[number_of_slide].text, self)
        self.print_text.setGeometry(300, 450, 700, 200)
        font = QtGui.QFont()
        font.setPointSize(16)
        self.print_text.setFont(font)

        self.create_new_slide.clicked.connect(self.new_slide)

        self.delete_button = QPushButton(self)
        self.delete_button.setText('Удалить этот слайд')
        self.delete_button.setStyleSheet('''QPushButton {            
                                    background: rgb(255,20,0);
                                    border-style: outset;
                                    border-width: px;
                                    border-radius: 15px;
                                    border-color: rgb(255,30,0);
                                    padding: 4px;
                                    color: rgb(255, 255, 255);
                                }''')
        font = QtGui.QFont()
        font.setPointSize(12)
        self.delete_button.setFont(font)
        self.delete_button.clicked.connect(self.delete_slide)

        self.pushButton.clicked.connect(self.create_presentation)

    def create_presentation(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_text(self.print_text.toPlainText())
        super().create_presentation()

    def paintEvent(self, event):
        self.border = QPainter()
        self.border.begin(self)
        self.draw(event, self.border)
        self.border.end()

    def draw(self, event, border):
        border.setBrush(QColor(184, 232, 176))
        border.setPen(QColor(184, 232, 176))
        border.drawRoundedRect(int(0.2336 * self.width()),
                               int(0.0650 * self.height()),
                               int(0.7300 * self.width()),
                               int(0.7950 * self.height()),
                               int(0.0475 * self.width()),
                               int(0.0475 * self.height()))

    def resizing2(self):
        self.print_title.setGeometry(int(0.2774 * self.width()),
                                     int(0.1200 * self.height()),
                                     int(0.6350 * self.width()),
                                     int(0.4100 * self.height()))
        self.print_text.setGeometry(int(0.2774 * self.width()),
                                    int(0.5760 * self.height()),
                                    int(0.6350 * self.width()),
                                    int(0.2400 * self.height()))
        self.delete_button.setGeometry(int(0.4659 * self.width()),
                                       int(0.9000 * self.height()),
                                       int(0.2500 * self.width()),
                                       int(0.0700 * self.height()))



    def resizeEvent(self, event):
        self.resized.emit()
        return super().resizeEvent(event)

    def new_slide(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_text(self.print_text.toPlainText())
        self.a = ChooseTemplate(self.slides, (self.width(), self.height()))
        self.a.show()
        self.hide()

    def to_slide(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_text(self.print_text.toPlainText())
        super().to_slide()


class Template2(MainScreen):
    resized = QtCore.pyqtSignal()

    def __init__(self, slides, number_of_slide, size):
        super().__init__(slides, size)
        self.number_of_slide = number_of_slide
        self.tempinit(number_of_slide)
        for i in self.list_of_buttons:
            i.clicked.connect(self.to_slide)

    def tempinit(self, number_of_slide):
        self.resized.connect(self.resizing3)
        self.print_title = QtWidgets.QPlainTextEdit(self.slides[number_of_slide].title, self)
        self.print_title.setGeometry(300, 150, 700,  150)
        font = QtGui.QFont()
        font.setPointSize(40)
        self.print_title.setFont(font)

        self.print_text = QtWidgets.QPlainTextEdit(self.slides[number_of_slide].text, self)
        self.print_text.setGeometry(300, 350, 700, 400)
        font = QtGui.QFont()
        font.setPointSize(16)
        self.print_text.setFont(font)

        self.create_new_slide.clicked.connect(self.new_slide)

        self.delete_button = QPushButton(self)
        self.delete_button.setText('Удалить этот слайд')
        self.delete_button.setGeometry(250, 835, 800, 30)
        self.delete_button.setStyleSheet('''QPushButton {            
                                            background: rgb(255,20,0);
                                            border-style: outset;
                                            border-width: px;
                                            border-radius: 15px;
                                            border-color: rgb(255,30,0);
                                            padding: 4px;
                                            color: rgb(255, 255, 255);
                                        }''')
        font = QtGui.QFont()
        font.setPointSize(12)
        self.delete_button.setFont(font)
        self.delete_button.clicked.connect(self.delete_slide)

        self.pushButton.clicked.connect(self.create_presentation)


    def create_presentation(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_text(self.print_text.toPlainText())
        super().create_presentation()

    def paintEvent(self, event):
        self.border = QPainter()
        self.border.begin(self)
        self.draw(event, self.border)
        self.border.end()

    def draw(self, event, border):
        border.setBrush(QColor(184, 232, 176))
        border.setPen(QColor(184, 232, 176))
        border.drawRoundedRect(int(0.2336 * self.width()),
                               int(0.0650 * self.height()),
                               int(0.7300 * self.width()),
                               int(0.7950 * self.height()),
                               int(0.0475 * self.width()),
                               int(0.0475 * self.height()))

    def new_slide(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_text(self.print_text.toPlainText())
        self.a = ChooseTemplate(self.slides, (self.width(), self.height()))
        self.a.show()
        self.hide()

    def to_slide(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_text(self.print_text.toPlainText())
        super().to_slide()

    def resizing3(self):
        self.print_title.setGeometry(int(0.2774 * self.width()),
                                     int(0.1200 * self.height()),
                                     int(0.6350 * self.width()),
                                     int(0.1800 * self.height()))
        self.print_text.setGeometry(int(0.2774 * self.width()),
                                    int(0.3590 * self.height()),
                                    int(0.6350 * self.width()),
                                    int(0.4550 * self.height()))
        self.delete_button.setGeometry(int(0.4659 * self.width()),
                                       int(0.9000 * self.height()),
                                       int(0.2500 * self.width()),
                                       int(0.0700 * self.height()))

    def resizeEvent3(self, event):
        self.resized.emit()
        return super().resizeEvent(event)


class Template3(MainScreen):
    resized = QtCore.pyqtSignal()

    def __init__(self, slides, number_of_slide, size):
        super().__init__(slides, size)
        self.number_of_slide = number_of_slide
        self.tempinit(number_of_slide)
        for i in self.list_of_buttons:
            i.clicked.connect(self.to_slide)

    def tempinit(self, number_of_slide):
        self.resized.connect(self.resizing4)
        self.print_title = QtWidgets.QPlainTextEdit(self.slides[number_of_slide].title, self)
        self.print_title.setGeometry(300, 150, 700,  150)
        font = QtGui.QFont()
        font.setPointSize(40)
        self.print_title.setFont(font)

        self.fname = self.slides[number_of_slide].picture
        self.input_picture = QPixmap(self.fname)

        self.print_picture = QtWidgets.QLabel(self)
        self.print_picture.setGeometry(300, 350, 700, 400)
        self.print_picture.setPixmap(self.input_picture)

        self.set_new_picture = QtWidgets.QPushButton(self)
        self.set_new_picture.setText('Вставить новую картинку')
        self.set_new_picture.setStyleSheet('''QPushButton {            
                                                    background: rgb(100,149,237);
                                                    border-style: outset;
                                                    border-width: px;
                                                    border-radius: 15px;
                                                    border-color: rgb(255,30,0);
                                                    padding: 4px;
                                                    color: rgb(255, 255, 255);
                                                }''')
        font = QtGui.QFont()
        font.setPointSize(10)
        self.set_new_picture.setFont(font)
        self.set_new_picture.clicked.connect(self.setting_new_picture)

        self.create_new_slide.clicked.connect(self.new_slide)

        self.delete_button = QPushButton(self)
        self.delete_button.setText('Удалить этот слайд')
        self.delete_button.setGeometry(250, 835, 800, 30)
        self.delete_button.setStyleSheet('''QPushButton {            
                                            background: rgb(255,20,0);
                                            border-style: outset;
                                            border-width: px;
                                            border-radius: 15px;
                                            border-color: rgb(255,30,0);
                                            padding: 4px;
                                            color: rgb(255, 255, 255);
                                        }''')
        font = QtGui.QFont()
        font.setPointSize(12)
        self.delete_button.setFont(font)
        self.delete_button.clicked.connect(self.delete_slide)

        self.pushButton.clicked.connect(self.create_presentation)

    def setting_new_picture(self):
        self.fname = QFileDialog.getOpenFileName(self, 'Выбрать картинку', '')[0]
        self.input_picture = QPixmap(self.fname)

        self.print_picture.setPixmap(self.input_picture)

    def create_presentation(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_picture(self.fname)
        super().create_presentation()

    def paintEvent(self, event):
        self.border = QPainter()
        self.border.begin(self)
        self.draw(event, self.border)
        self.border.end()

    def draw(self, event, border):
        border.setBrush(QColor(184, 232, 176))
        border.setPen(QColor(184, 232, 176))
        border.drawRoundedRect(int(0.2336 * self.width()),
                               int(0.0650 * self.height()),
                               int(0.7300 * self.width()),
                               int(0.7950 * self.height()),
                               int(0.0475 * self.width()),
                               int(0.0475 * self.height()))

    def new_slide(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_picture(self.fname)
        self.a = ChooseTemplate(self.slides, (self.width(), self.height()))
        self.a.show()
        self.hide()

    def to_slide(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_picture(self.fname)
        super().to_slide()

    def resizing4(self):
        self.print_title.setGeometry(int(0.2774 * self.width()),
                                     int(0.1200 * self.height()),
                                     int(0.6350 * self.width()),
                                     int(0.1800 * self.height()))
        self.print_picture.setGeometry(int(0.2774 * self.width()),
                                       int(0.3590 * self.height()),
                                       int(0.6350 * self.width()),
                                       int(0.4550 * self.height()))

        self.delete_button.setGeometry(int(0.4659 * self.width()),
                                       int(0.9000 * self.height()),
                                       int(0.2500 * self.width()),
                                       int(0.0700 * self.height()))
        self.set_new_picture.setGeometry(int(0.2336 * self.width()),
                                         int(0.9000 * self.height()),
                                         int(0.2250 * self.width()),
                                         int(0.0700 * self.height()))

    def resizeEvent4(self, event):
        self.resized.emit()
        return super().resizeEvent(event)


class Template4(MainScreen):
    resized = QtCore.pyqtSignal()

    def __init__(self, slides, number_of_slide, size):
        super().__init__(slides, size)
        self.number_of_slide = number_of_slide
        self.tempinit(number_of_slide)
        for i in self.list_of_buttons:
            i.clicked.connect(self.to_slide)

    def tempinit(self, number_of_slide):
        self.resized.connect(self.resizing5)
        self.print_title = QtWidgets.QPlainTextEdit(self.slides[number_of_slide].title, self)
        self.print_title.setGeometry(300, 150, 700,  150)
        font = QtGui.QFont()
        font.setPointSize(40)
        self.print_title.setFont(font)

        self.fname = self.slides[number_of_slide].picture
        self.input_picture = QPixmap(self.fname)

        self.print_picture = QtWidgets.QLabel(self)
        self.print_picture.setGeometry(300, 350, 350, 400)
        self.print_picture.setPixmap(self.input_picture)

        self.print_text = QtWidgets.QPlainTextEdit(self.slides[number_of_slide].text, self)
        self.print_text.setGeometry(650, 350, 350, 400)
        font = QtGui.QFont()
        font.setPointSize(16)
        self.print_text.setFont(font)

        self.set_new_picture = QtWidgets.QPushButton(self)
        self.set_new_picture.setText('Вставить новую картинку')
        self.set_new_picture.setStyleSheet('''QPushButton {            
                                                            background: rgb(100,149,237);
                                                            border-style: outset;
                                                            border-width: px;
                                                            border-radius: 15px;
                                                            border-color: rgb(255,30,0);
                                                            padding: 4px;
                                                            color: rgb(255, 255, 255);
                                                        }''')
        font = QtGui.QFont()
        font.setPointSize(10)
        self.set_new_picture.setFont(font)
        self.set_new_picture.clicked.connect(self.setting_new_picture)

        self.create_new_slide.clicked.connect(self.new_slide)

        self.delete_button = QPushButton(self)
        self.delete_button.setText('Удалить этот слайд')
        self.delete_button.setGeometry(250, 835, 800, 30)
        self.delete_button.setStyleSheet('''QPushButton {            
                                            background: rgb(255,20,0);
                                            border-style: outset;
                                            border-width: px;
                                            border-radius: 15px;
                                            border-color: rgb(255,30,0);
                                            padding: 4px;
                                            color: rgb(255, 255, 255);
                                        }''')
        font = QtGui.QFont()
        font.setPointSize(12)
        self.delete_button.setFont(font)
        self.delete_button.clicked.connect(self.delete_slide)

        self.pushButton.clicked.connect(self.create_presentation)

    def setting_new_picture(self):
        self.fname = QFileDialog.getOpenFileName(self, 'Выбрать картинку', '')[0]
        self.input_picture = QPixmap(self.fname)
        self.print_picture.setPixmap(self.input_picture)

    def create_presentation(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_picture(self.fname)
        self.slides[self.number_of_slide].set_text(self.print_text.toPlainText())
        super().create_presentation()

    def paintEvent(self, event):
        self.border = QPainter()
        self.border.begin(self)
        self.draw(event, self.border)
        self.border.end()

    def draw(self, event, border):
        border.setBrush(QColor(184, 232, 176))
        border.setPen(QColor(184, 232, 176))
        border.drawRoundedRect(int(0.2336 * self.width()),
                               int(0.0650 * self.height()),
                               int(0.7300 * self.width()),
                               int(0.7950 * self.height()),
                               int(0.0475 * self.width()),
                               int(0.0475 * self.height()))

    def new_slide(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_picture(self.fname)
        self.slides[self.number_of_slide].set_text(self.print_text.toPlainText())
        self.a = ChooseTemplate(self.slides, (self.width(), self.height()))
        self.a.show()
        self.hide()

    def to_slide(self):
        self.slides[self.number_of_slide].set_title(self.print_title.toPlainText())
        self.slides[self.number_of_slide].set_picture(self.fname)
        self.slides[self.number_of_slide].set_text(self.print_text.toPlainText())
        super().to_slide()

    def resizing5(self):
        self.print_title.setGeometry(int(0.2774 * self.width()),
                                     int(0.1200 * self.height()),
                                     int(0.6350 * self.width()),
                                     int(0.1800 * self.height()))
        self.print_picture.setGeometry(int(0.2774 * self.width()),
                                       int(0.3590 * self.height()),
                                       int(0.3175 * self.width()),
                                       int(0.4550 * self.height()))
        self.print_text.setGeometry(int(0.5909 * self.width()),
                                    int(0.3590 * self.height()),
                                    int(0.3175 * self.width()),
                                    int(0.4550 * self.height()))
        self.delete_button.setGeometry(int(0.4659 * self.width()),
                                       int(0.9000 * self.height()),
                                       int(0.2500 * self.width()),
                                       int(0.0700 * self.height()))
        self.set_new_picture.setGeometry(int(0.2336 * self.width()),
                                         int(0.9000 * self.height()),
                                         int(0.2250 * self.width()),
                                         int(0.0700 * self.height()))

    def resizeEvent5(self, event):
        self.resized.emit()
        return super().resizeEvent(event)


if __name__ == '__main__':
    app = QtWidgets.QApplication(sys.argv)
    ex = GreetingWindow()
    ex.show()
    sys.exit(app.exec())
